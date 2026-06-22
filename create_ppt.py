import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# 16:9 Aspect Ratio widescreen format
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Visual design theme color constants (Matching Liquid Glass System)
DARK_BG = RGBColor(5, 7, 12)
TEXT_WHITE = RGBColor(255, 255, 255)
TEAL_MOMENTUM = RGBColor(119, 242, 228)
SKY_BLUE_FLOW = RGBColor(107, 182, 243)
SUCCESS_EMERALD = RGBColor(16, 185, 129)
PROTECTIVE_ROSE = RGBColor(239, 68, 68)
SLATE_GRAY = RGBColor(148, 163, 184)
SLATE_CARD_BG = RGBColor(13, 20, 35) # Translucent-looking deep slate backing
SLATE_BORDER = RGBColor(45, 58, 78) # Crisp glass highlights border

def apply_slide_base(slide, title_text):
    # Set dark background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background() # Hide borders
    
    # BACKGROUND CAUSTICS (Physically simulates ambient colorful lights blurred behind frosted glass panels!)
    # Orb 1 (Sky Blue)
    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.2), Inches(3.5), Inches(3.5))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = RGBColor(10, 32, 54) # Sky Blue caustics glow
    orb1.line.fill.background()
    
    # Orb 2 (Teal)
    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(3.8), Inches(3.2), Inches(3.2))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = RGBColor(8, 38, 42) # Teal caustics glow
    orb2.line.fill.background()
    
    # Elegant Top Header bar backing (frosted layout layer)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = SLATE_CARD_BG
    hdr.line.color.rgb = SLATE_BORDER
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.0), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL_MOMENTUM
    
    # Meta / Brand Tracker Text Box
    txBox_track = slide.shapes.add_textbox(Inches(8.8), Inches(0.45), Inches(3.8), Inches(0.7))
    tf_track = txBox_track.text_frame
    tf_track.word_wrap = True
    p_track = tf_track.paragraphs[0]
    p_track.alignment = PP_ALIGN.RIGHT
    p_track.text = "WayaX Frosted Glass Spec • INH00010876"
    p_track.font.name = 'Calibri'
    p_track.font.size = Pt(10)
    p_track.font.color.rgb = SLATE_GRAY

def create_info_card(slide, x, y, width, height, title, lines):
    # Frosted Panel Backing card (simulates frosted glass with custom borders)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = SLATE_CARD_BG
    card.line.color.rgb = SLATE_BORDER # Crisp frosted glass border
    
    # Title
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), width - Inches(0.3), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEAL_MOMENTUM # Vibrant theme highlight
    
    # Content Paragraphs
    txBox_body = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.65), width - Inches(0.3), height - Inches(0.8))
    tf_body = txBox_body.text_frame
    tf_body.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p_body = tf_body.paragraphs[0]
        else:
            p_body = tf_body.add_paragraph()
        p_body.text = line
        p_body.font.name = 'Calibri'
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = TEXT_WHITE
        p_body.space_after = Pt(6)

# ==========================================
# SLIDE 1: COVER WITH HIGH-FIDELITY FROSTED REFRACTIVE GLOWS
# ==========================================
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Base background cover
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# Background ambient caustics glows (Teal & Sky Blue underlays)
cover_orb1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(1.2), Inches(5.8), Inches(4.8))
cover_orb1.fill.solid()
cover_orb1.fill.fore_color.rgb = RGBColor(12, 38, 54) # Deep Sky Blue glow
cover_orb1.line.fill.background()

cover_orb2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.8), Inches(5.2), Inches(4.2))
cover_orb2.fill.solid()
cover_orb2.fill.fore_color.rgb = RGBColor(10, 42, 45) # Deep Teal glow
cover_orb2.line.fill.background()

# Central Frosted Glass Panel Card
glass_panel = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.8), Inches(11.93), Inches(3.8))
glass_panel.fill.solid()
glass_panel.fill.fore_color.rgb = SLATE_CARD_BG
glass_panel.line.color.rgb = SLATE_BORDER

# Title text inside the glass card
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

# Main title
p1 = tf.paragraphs[0]
p1.text = "WayaX Liquid Glass Design System"
p1.font.name = 'Arial'
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEAL_MOMENTUM
p1.space_after = Pt(8)

# Subtitle
p2 = tf.add_paragraph()
p2.text = "Visual Specifications, Specular Refractions & SEBI Advisory Handoff Rules"
p2.font.name = 'Calibri'
p2.font.size = Pt(18)
p2.font.color.rgb = SKY_BLUE_FLOW
p2.space_after = Pt(20)

# Auditor & Model info
p3 = tf.add_paragraph()
p3.text = "SEBI Registered Research Analyst: INH00010876 | Programmatically Derived via Qwen2.5-Coder:1.5B"
p3.font.name = 'Calibri'
p3.font.size = Pt(12)
p3.font.italic = True
p3.font.color.rgb = SLATE_GRAY

# Specular accent highlight strip
strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.2), Inches(4.5), Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = TEAL_MOMENTUM
strip.line.fill.background()

# ==========================================
# SLIDE 2: CORE PHILOSOPHY
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
apply_slide_base(slide2, "CORE PHILOSOPHY: REFRACTIVE SPECULARITY")

# Left Column Big Callout
callout = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.8))
tf_call = callout.text_frame
tf_call.word_wrap = True
p_call = tf_call.paragraphs[0]
p_call.text = "Rejecting Flat, Synthetic UI Cards"
p_call.font.name = 'Arial'
p_call.font.size = Pt(28)
p_call.font.bold = True
p_call.font.color.rgb = TEAL_MOMENTUM
p_call.space_after = Pt(16)

p_desc = tf_call.add_paragraph()
p_desc.text = "WayaX mimics physical optical glass layered over responsive grids to establish deep visceral trust in audited advisory data. The interface feels organic, structural, and visually alive."
p_desc.font.name = 'Calibri'
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = TEXT_WHITE
p_desc.space_after = Pt(12)

# Right Column - 3 Pillars Cards
create_info_card(slide2, Inches(5.5), Inches(1.8), Inches(7.3), Inches(1.4), 
                 "1. Optical Depth Separation", 
                 ["Physical layered dimensions divide user conversational bubbles from primary, audited financial recommendations.", "Provides clear hierarchy and prevents information fatigue."])

create_info_card(slide2, Inches(5.5), Inches(3.4), Inches(7.3), Inches(1.4), 
                 "2. Backdrop Caustics Refractions", 
                 ["Accent gradient colors and ambient refractions pass underneath card panels to simulate natural light behavior.", "Enhances interactivity and visual feedback during scroll states."])

create_info_card(slide2, Inches(5.5), Inches(5.0), Inches(7.3), Inches(1.4), 
                 "3. Sub-Pixel Specular Highlights", 
                 ["Top borders are detailed with a 1px specular white highlight at 12% opacity to emulate precise glass beveling.", "Ensures crisp contours against dynamic background details."])

# ==========================================
# SLIDE 3: FOUNDATIONS & DESIGN TOKENS
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
apply_slide_base(slide3, "FOUNDATIONS & BRAND VALUE TOKENS")

# Colors Card
create_info_card(slide3, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), 
                 "A. Liquid Specular Color Scale", 
                 [
                     "• TEAL MOMENTUM (#77F2E4): Left logo gradient stop (0%). Highlights core brand actions.",
                     "• SKY BLUE FLOW (#6BB6F3): Right logo gradient stop (100%). Soft link primary background accent.",
                     "• SUCCESS EMERALD (#10B981): Highlights positive growth indices, BUY calls, target prices.",
                     "• PROTECTIVE ROSE (#EF4444): Stop loss limits, SHORT triggers, and excessive leverage warnings.",
                     "• CANVAS BLACK (#000000) & SLATE BACKING (#0D0E12): High contrast dark environment layout.",
                     "• DEPRECATION: Salted Camel (#D5AD76) is fully deprecated to enforce maximum design clarity."
                 ])

# Typography Card
create_info_card(slide3, Inches(6.8), Inches(1.8), Inches(6.0), Inches(2.3), 
                 "B. Typography System Scales", 
                 [
                     "• Sora Display font curve scales for titles and onboarding structures (28px Bold / 20px / 17px).",
                     "• Inter Body sans-serif scales for streamed response advisory paragraphs (13px Medium).",
                     "• JetBrains Mono data text scales for clean metrics, prices, and ratios readability (12.5px Bold)."
                 ])

# Grid Spacing Card
create_info_card(slide3, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.3), 
                 "C. Spacing Grid laws (8px Base Scale)", 
                 [
                     "• px-1 (4px): Micro status indicator pulses, margins.",
                     "• px-1.5 (6px): Advisory monospace symbol badge margins.",
                     "• p-2.5 (10px): Search bar icon alignments, categories gaps.",
                     "• p-4.5 (18px): Standard table row vertical padding boundaries.",
                     "• p-6 (24px): Spacious container gaps on active screens."
                 ])

# ==========================================
# SLIDE 4: INTERACTIVE COMPONENTS
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
apply_slide_base(slide4, "INTERACTIVE COMPONENTS SPECIFICATION")

# Buttons card
create_info_card(slide4, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), 
                 "A. Interactive Solid Button Matrix", 
                 [
                     "1. Primary Action Buttons (Capsule):",
                     "   - Default: solid bg-indigo-600, Sora Bold, rounded-2xl.",
                     "   - Hover: bg-indigo-50, translateY(-1.5px) for light elevation.",
                     "   - Active: Scale down to 0.95, bg-indigo-750.",
                     "   - Disabled: opacity-50, cursor-not-allowed, bg-slate-800.",
                     "2. Secondary Glass Buttons (.liquid-glass-button):",
                     "   - Default: bg-white/[0.05], 1px white border at 10% opacity.",
                     "   - Hover: Triggers diagonal specular specular sweeper at 38deg.",
                     "   - Active: Scale 0.98, background shifts to bg-white/[0.12]."
                 ])

# Search input card
create_info_card(slide4, Inches(6.8), Inches(1.8), Inches(6.0), Inches(2.3), 
                 "B. Recessed Search Input Well", 
                 [
                     "• Default (.liquid-glass-input): bg-black/25, inner shadow inset shadow [inset_0_2px_4px_rgba(0,0,0,0.35)].",
                     "• Focus outline: border-white/20, bg-black/35, releases a beautiful white outline breathing glow to highlight focus."
                 ])

# Dropdown limits card
create_info_card(slide4, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.3), 
                 "C. Presets Selector Panel Constraint", 
                 [
                     "• Height constraint: locked to max-h-[132px] (mobile) or max-h-[180px] (desktop).",
                     "• Visual Law: Displays exactly 3 questions + a 30% visual vertical cut-off hint of the 4th item.",
                     "• Design Value: Eliminates layout shift scroll jumps when expanding dropdown panels."
                 ])

# ==========================================
# SLIDE 5: PATTERN LAWS & DEV HANDOFF
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
apply_slide_base(slide5, "PATTERN LAWS & DEV HANDOFF COMPLIANCE")

# Emerald Law card
create_info_card(slide5, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), 
                 "A. The Emerald Green Law (BUY Actions)", 
                 [
                     "• Renders strictly in Success Emerald (#10B981) to emphasize positive growth parameters.",
                     "• BUY call badges: semi-translucent green capsule (bg-emerald-500/10, border-emerald-500/20) with a pulsing green indicator dot (animate-pulse).",
                     "• Target Upside badges: grouped with TrendingUp icon and bold monospace green text.",
                     "• Positive return metrics (1M, 3M, 1Y period growths) must highlight green."
                 ])

# Rose Law card
create_info_card(slide5, Inches(6.8), Inches(1.8), Inches(6.0), Inches(2.3), 
                 "B. The Rose Red Law (Protection & Warnings)", 
                 [
                     "• Renders strictly in Protective Rose (#EF4444) to ensure immediate protective visibility.",
                     "• SHORT call badges: semi-translucent red capsule with pulsing red indicator dot.",
                     "• Protective Stop Loss (SL): highlighted in bold monospace red text.",
                     "• Volatility indices and negative period periods must highlight red."
                 ])

# Do's and Dont's card
create_info_card(slide5, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.3), 
                 "C. Critical Handoff Guidelines", 
                 [
                     "• DO: Pad all row indices with leading zeros (e.g. 01, 02) and align centered.",
                     "• DO: Snap scroll alignment to top of streamed chat responses to prevent scroll jumping.",
                     "• DON'T: Inject deprecated Salted Camel (#D5AD76) inside card color stops.",
                     "• DON'T: Render scrollbars without the translucent glassmorphic thumbs."
                 ])

# ==========================================
# SLIDE 6: DATA ORGANISMS & SYNC CHANNELS
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
apply_slide_base(slide6, "STOCK TABLE ORGANISM & FIGMA SYNC")

# Stock Table Specs Card
create_info_card(slide6, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), 
                 "A. SEBI-Compliant Stock Table Specs", 
                 [
                     "• Column Width margins strictly defined to prevent layout shift during row accordion expansion:",
                     "  - Index: w-14, centered, monospaced.",
                     "  - Company / Symbol: min-w-[220px], left-aligned.",
                     "  - Call: w-32, status badging.",
                     "  - Numeric Metrics: center-aligned monospaced parameters.",
                     "• Hover transitions: py-4.5 (18px) padding, 300ms hover bg-white/[0.035].",
                     "• Accordion drawer divisions: divided into Technical Matrix column, Fundamental Shield column, and Investment Thesis signature."
                 ])

# Figma Synchronization Card
create_info_card(slide6, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8), 
                 "B. Figma & Code Synchronization Channel", 
                 [
                     "• 1. Figma Drag-and-Drop Vector Spec Sheet:",
                     "  - Source: C:\\Users\\lcppd\\.gemini\\antigravity\\scratch\\WayaX\\WayaX_Design_System.svg",
                     "  - Drag-and-drop this file directly into the Figma canvas to instantly import all 6 visual pages as fully editable vector shapes and layers.",
                     "• 2. CSS Utility Synchronization:",
                     "  - Source: C:\\Users\\lcppd\\.gemini\\antigravity\\scratch\\WayaX\\src\\index.css",
                     "  - Contains all .liquid-glass class details.",
                     "• 3. React implementation:",
                     "  - Source: C:\\Users\\lcppd\\.gemini\\antigravity\\scratch\\WayaX\\src\\App.tsx",
                     "  - Maps viewport spacings, tokens, and margins."
                 ])

# Save presentation
prs.save("WayaX_Frosted_Glass_Design_System.pptx")
print("PowerPoint presentation with frosted glassmorphic styles generated successfully!")
